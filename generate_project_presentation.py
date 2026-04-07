from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT_FILE = "AI_Resume_Analyzer_Project_Presentation.pptx"


def set_background(slide, color=(245, 248, 255)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 40, 90)


def add_subtitle(slide, text):
    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.0), Inches(0.7))
    tf = sub_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(70, 85, 120)


def add_bullets(slide, bullets, left=0.8, top=1.9, width=11.6, height=5.0, font_size=22):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for idx, line in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line, tuple):
            text, level = line
        else:
            text, level = line, 0
        p.text = text
        p.level = level
        p.font.size = Pt(font_size - (2 if level > 0 else 0))
        p.font.color.rgb = RGBColor(35, 45, 65)
        p.space_after = Pt(8)


def add_footer(slide, text="AI Resume & Job Match Analyzer | Emerging Tech Project"):
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12.0), Inches(0.35))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(100, 110, 135)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # Slide 1: Title with requested details
    slide = prs.slides.add_slide(blank)
    set_background(slide, (236, 244, 255))
    add_title(slide, "AI Resume & Job Match Analyzer")
    add_subtitle(slide, "Project Presentation")

    details = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.8), Inches(3.2))
    tf = details.text_frame
    tf.clear()
    lines = [
        "Name: Vaibhav Krishali",
        "Registeration Number: 23FE10CSE00301",
        "Section: F",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = RGBColor(18, 39, 89)
        p.space_after = Pt(16)

    add_footer(slide)

    # Slide 2: Problem and objective
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Problem Statement and Objective")
    add_bullets(
        slide,
        [
            "Recruiters and ATS systems filter resumes quickly, often rejecting strong candidates for poor keyword alignment.",
            "Manual resume tailoring for every job role is time-consuming and inconsistent.",
            "Objective: Build an AI + ML system that scans resumes, matches them to target roles/JDs, and suggests concrete improvements.",
            "Outcome: Faster, role-focused resume optimization with structured, actionable feedback.",
        ],
    )
    add_footer(slide)

    # Slide 3: What the project does
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Core Features")
    add_bullets(
        slide,
        [
            "Resume input via PDF upload and text extraction.",
            "CV-first scan for ATS readiness score and profile status.",
            "Role suggestion engine based on candidate profile.",
            "Job-match analysis with matching skills, missing skills, interview questions, and ATS suggestions.",
            "Resume upgrade flow with before vs after score and rewritten bullet points.",
            "Overleaf-ready LaTeX resume generation, plus PDF and preview download.",
            "Dedicated ML endpoints using scikit-learn for quality scoring and semantic matching.",
        ],
        font_size=21,
    )
    add_footer(slide)

    # Slide 4: Architecture
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "System Architecture")
    add_bullets(
        slide,
        [
            "Frontend: Single-page UI built with HTML, CSS, and vanilla JavaScript.",
            "Backend: FastAPI app with clear route separation (web routes and API routes).",
            "Services layer handles business logic: analyzer orchestration, PDF parsing, LLM connectors, LaTeX compile, preview rendering.",
            "Data flow:",
            ("Client -> API -> ResumeJobAnalyzer -> (OpenAI / Groq / Local LLM / Free fallback)", 1),
            ("Client -> API -> ML analyzer (TF-IDF + Logistic Regression + Cosine Similarity)", 1),
            ("Client -> API -> Resume template + LaTeX compiler -> PDF/PNG preview", 1),
        ],
        font_size=20,
    )
    add_footer(slide)

    # Slide 5: Tech stack
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Technology Stack")
    add_bullets(
        slide,
        [
            "Backend Framework: FastAPI + Uvicorn",
            "Frontend: Jinja templates + Vanilla JS + CSS",
            "NLP/ML: scikit-learn, NumPy",
            "Modeling methods:",
            ("TF-IDF vectorization (500 features, 1-2 grams)", 1),
            ("Logistic Regression classifier for quality score", 1),
            ("Cosine Similarity for resume-job semantic fit", 1),
            "Integrations: OpenAI API, Groq API, local Ollama LLM",
            "Document pipeline: pypdf parsing, LaTeX generation/compilation",
            "Deployment: Docker support",
        ],
        font_size=20,
    )
    add_footer(slide)

    # Slide 6: API design
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Major API Endpoints")
    add_bullets(
        slide,
        [
            "POST /api/scan-cv: CV score, role direction, strengths, improvement areas.",
            "POST /api/suggest-roles: suggested target roles from resume content.",
            "POST /api/analyze: full role/JD-based ATS analysis and recommendations.",
            "POST /api/resume-upgrade: role-specific improved draft, score delta, and LaTeX output.",
            "POST /api/compile-latex and /api/compile-latex-preview-image: compile and render output resume.",
            "ML APIs: /api/ml-insights, /api/ml-quality-score, /api/ml-skill-similarity.",
        ],
        font_size=20,
    )
    add_footer(slide)

    # Slide 7: ML pipeline
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Machine Learning Pipeline")
    add_bullets(
        slide,
        [
            "Step 1: Convert resume and job text to TF-IDF vectors.",
            "Step 2: Predict resume quality using trained Logistic Regression classifier.",
            "Step 3: Compute resume-job cosine similarity for semantic fit.",
            "Step 4: Extract NLP features (word count, vocabulary richness, technical keyword density).",
            "Step 5: Generate interpretable output: score, confidence, matched/missing skills, and feature metrics.",
            "Performance: warm inference under ~50ms for most requests after model load.",
        ],
        font_size=20,
    )
    add_footer(slide)

    # Slide 8: User workflow
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "End-to-End User Workflow")
    add_bullets(
        slide,
        [
            "1) Upload PDF resume.",
            "2) Run CV Scan for baseline ATS score and role suggestions.",
            "3) Choose target role and optionally paste job description.",
            "4) Run analysis to receive skill gap, ATS improvements, and interview questions.",
            "5) Generate upgraded resume and compare old vs new snapshot.",
            "6) Export final resume via LaTeX/PDF output.",
        ],
        font_size=22,
    )
    add_footer(slide)

    # Slide 9: Strengths
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Project Strengths")
    add_bullets(
        slide,
        [
            "Multi-mode AI engine with robust fallback strategy (OpenAI -> Groq -> local LLM -> heuristic).",
            "Combines explainable ML scoring with generative AI suggestions.",
            "Practical resume engineering outputs: rewritten bullets, ATS keywords, tailored summaries.",
            "Production-friendly structure: modular services, schema-driven responses, Docker support.",
            "Useful for students and job seekers who need rapid role-specific resume adaptation.",
        ],
        font_size=21,
    )
    add_footer(slide)

    # Slide 10: Limitations and future scope
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Limitations and Future Enhancements")
    add_bullets(
        slide,
        [
            "Current ML training set is synthetic and relatively small.",
            "Job-domain ontology can be expanded for more fine-grained role matching.",
            "Potential improvements:",
            ("Larger curated dataset for model training and benchmark evaluation", 1),
            ("RAG-based JD/resume knowledge grounding", 1),
            ("User accounts with history and version tracking", 1),
            ("Advanced analytics dashboard and interview simulation", 1),
        ],
        font_size=20,
    )
    add_footer(slide)

    # Slide 11: Conclusion
    slide = prs.slides.add_slide(blank)
    set_background(slide, (236, 244, 255))
    add_title(slide, "Conclusion")
    add_bullets(
        slide,
        [
            "This project delivers a complete AI + ML resume intelligence platform.",
            "It transforms raw resume input into measurable, role-oriented recommendations.",
            "The architecture is modular, extensible, and deployment-ready for real-world usage.",
            "Thank you.",
        ],
        top=2.2,
        font_size=26,
    )
    add_footer(slide, "Prepared by Vaibhav Krishali")

    prs.save(OUTPUT_FILE)


if __name__ == "__main__":
    build_presentation()
    print(f"Presentation generated: {OUTPUT_FILE}")
